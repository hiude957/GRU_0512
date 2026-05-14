#!/usr/bin/env python3
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_PATH = Path("docs/experiments/gru_msfix_autoregressive_report_20260513.pptx")

FONT = "Microsoft YaHei"
BG = RGBColor(246, 248, 251)
NAVY = RGBColor(22, 35, 57)
BLUE = RGBColor(43, 105, 179)
TEAL = RGBColor(26, 132, 126)
ORANGE = RGBColor(222, 124, 46)
GREEN = RGBColor(68, 145, 91)
RED = RGBColor(188, 82, 82)
GRAY = RGBColor(101, 113, 130)
LIGHT_BLUE = RGBColor(225, 236, 250)
LIGHT_TEAL = RGBColor(222, 244, 241)
LIGHT_ORANGE = RGBColor(255, 238, 221)
WHITE = RGBColor(255, 255, 255)


def set_slide_bg(slide, color: RGBColor = BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def set_run_font(run, size: int, color: RGBColor = NAVY, bold: bool = False) -> None:
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def add_textbox(slide, x, y, w, h, text: str, size: int = 18, color: RGBColor = NAVY, bold: bool = False,
                align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run_font(run, size=size, color=color, bold=bold)
    return box


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    add_textbox(slide, Inches(0.55), Inches(0.35), Inches(11.9), Inches(0.45), title, 26, NAVY, True)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.95), Inches(1.1), Inches(0.04))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()
    if subtitle:
        add_textbox(slide, Inches(0.55), Inches(1.05), Inches(11.8), Inches(0.35), subtitle, 12, GRAY)


def add_footer(slide, page: int) -> None:
    add_textbox(slide, Inches(0.55), Inches(7.08), Inches(8), Inches(0.22),
                "GRU_0512 | livedata msfix + autoregressive validation", 8, GRAY)
    add_textbox(slide, Inches(12.15), Inches(7.08), Inches(0.6), Inches(0.22), str(page), 8, GRAY,
                align=PP_ALIGN.RIGHT)


def add_card(slide, x, y, w, h, title: str, body: list[str], accent: RGBColor = BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(220, 226, 235)
    shape.line.width = Pt(0.8)
    add_textbox(slide, x + Inches(0.22), y + Inches(0.16), w - Inches(0.44), Inches(0.28),
                title, 14, accent, True)
    text = "\n".join(f"• {item}" for item in body)
    add_textbox(slide, x + Inches(0.22), y + Inches(0.58), w - Inches(0.44), h - Inches(0.72),
                text, 11, NAVY)
    return shape


def add_metric(slide, x, y, w, h, label: str, value: str, note: str = "", color: RGBColor = BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(220, 226, 235)
    add_textbox(slide, x + Inches(0.12), y + Inches(0.12), w - Inches(0.24), Inches(0.25), label, 10, GRAY)
    add_textbox(slide, x + Inches(0.12), y + Inches(0.42), w - Inches(0.24), Inches(0.42), value, 22, color, True)
    if note:
        add_textbox(slide, x + Inches(0.12), y + Inches(0.90), w - Inches(0.24), Inches(0.28), note, 8, GRAY)


def add_table(slide, x, y, w, h, headers: list[str], rows: list[list[str]], font_size: int = 9,
              header_fill: RGBColor = NAVY):
    table = slide.shapes.add_table(len(rows) + 1, len(headers), x, y, w, h).table
    for col_idx in range(len(headers)):
        table.columns[col_idx].width = int(w / len(headers))
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        cell.text = header
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                set_run_font(run, font_size, WHITE, True)
    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if row_idx % 2 else RGBColor(241, 245, 250)
            cell.text = value
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if col_idx else PP_ALIGN.LEFT
                for run in p.runs:
                    set_run_font(run, font_size, NAVY, False)
    return table


def add_arrow(slide, x, y, w, h, color: RGBColor = GRAY) -> None:
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()


def add_process_box(slide, x, y, w, h, title: str, subtitle: str, fill: RGBColor):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = RGBColor(210, 218, 230)
    add_textbox(slide, x + Inches(0.12), y + Inches(0.13), w - Inches(0.24), Inches(0.28), title, 13, NAVY, True,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.12), y + Inches(0.48), w - Inches(0.24), h - Inches(0.55), subtitle, 9, GRAY,
                align=PP_ALIGN.CENTER)


def add_bar_chart(slide, x, y, w, h, labels: list[str], values: list[float], max_value: float,
                  title: str, color: RGBColor = BLUE) -> None:
    add_textbox(slide, x, y, w, Inches(0.28), title, 12, NAVY, True)
    plot_x = x + Inches(0.2)
    plot_y = y + Inches(0.55)
    bar_gap = Inches(0.12)
    bar_w = (w - Inches(0.4) - bar_gap * (len(labels) - 1)) / len(labels)
    for idx, (label, value) in enumerate(zip(labels, values)):
        bx = plot_x + idx * (bar_w + bar_gap)
        bh = (h - Inches(1.2)) * (value / max_value)
        by = plot_y + (h - Inches(1.2)) - bh
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx, by, bar_w, bh)
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        add_textbox(slide, bx - Inches(0.03), by - Inches(0.25), bar_w + Inches(0.06), Inches(0.22),
                    f"{value:.3f}", 8, NAVY, True, align=PP_ALIGN.CENTER)
        add_textbox(slide, bx - Inches(0.05), y + h - Inches(0.45), bar_w + Inches(0.1), Inches(0.25),
                    label, 8, GRAY, align=PP_ALIGN.CENTER)


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1. Title
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, NAVY)
    add_textbox(slide, Inches(0.8), Inches(1.15), Inches(11.8), Inches(0.65),
                "GRU 传感器预测与全天自回归验证汇报", 30, WHITE, True)
    add_textbox(slide, Inches(0.82), Inches(2.0), Inches(11.2), Inches(0.45),
                "基于 GRU_0512 no-mask cache、livedata 毫秒修复、20260508-20260511 验证集", 16,
                RGBColor(210, 222, 238))
    add_metric(slide, Inches(0.85), Inches(3.25), Inches(2.4), Inches(1.3), "best checkpoint", "epoch 7", "val_mae=0.006515", TEAL)
    add_metric(slide, Inches(3.55), Inches(3.25), Inches(2.4), Inches(1.3), "val-only rows", "2.136M", "5/8-5/11", BLUE)
    add_metric(slide, Inches(6.25), Inches(3.25), Inches(2.4), Inches(1.3), "daily closed-loop", "93.12%", "日内不重置", GREEN)
    add_metric(slide, Inches(8.95), Inches(3.25), Inches(2.4), Inches(1.3), "4-day chain", "92.41%", "跨天不重置", ORANGE)
    add_textbox(slide, Inches(0.85), Inches(6.55), Inches(5), Inches(0.3), "2026-05-13", 10, RGBColor(210, 222, 238))

    # 2. Objective
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "汇报目标", "本次汇报回答三个问题：模型为什么选 GRU、模型如何工作、全天自回归结果是否可用。")
    add_card(slide, Inches(0.75), Inches(1.55), Inches(3.8), Inches(4.7), "1. 数据与对齐", [
        "APC / livedata / log action 对齐到 110ms 网格",
        "只在 livedata 覆盖范围输出训练/验证行",
        "覆盖外 action 先更新 state，再从输出行中过滤",
        "修复 livedata 1/2 位毫秒解析问题",
    ], TEAL)
    add_card(slide, Inches(4.75), Inches(1.55), Inches(3.8), Inches(4.7), "2. GRU 训练", [
        "输入 401 维，不包含 mask 位",
        "输出 150 个 sensor 绝对值",
        "2 层 GRU，hidden_size=256",
        "epoch 7 达到 best validation MAE",
    ], BLUE)
    add_card(slide, Inches(8.75), Inches(1.55), Inches(3.8), Inches(4.7), "3. 全天自回归验证", [
        "单日全天：每天 warmup 后整天闭环",
        "四天连续：5/8 开头 warmup 一次，跨天继承",
        "比较重置策略对长链误差漂移的影响",
        "重点观察 continuous sensor 的长链稳定性",
    ], ORANGE)
    add_footer(slide, 2)

    # 3. Why GRU
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "为什么选择 GRU", "当前问题是高频时序预测，并且需要流式继承状态。GRU 是复杂度和效果之间的稳健选择。")
    add_card(slide, Inches(0.75), Inches(1.45), Inches(5.8), Inches(2.15), "数据形态匹配", [
        "110ms 固定网格形成强时序依赖",
        "sensor 值、evt/action、state_* 都随时间连续演化",
        "预测时需要把上一时刻 sensor 预测回灌到下一步输入",
    ], TEAL)
    add_card(slide, Inches(6.8), Inches(1.45), Inches(5.75), Inches(2.15), "状态继承能力", [
        "GRU hidden 可以自然承载设备过程上下文",
        "适合日内 segment 边界和跨天边界的连续继承",
        "推理时只需保存 hidden state，部署逻辑简单",
    ], BLUE)
    add_card(slide, Inches(0.75), Inches(3.85), Inches(5.8), Inches(2.15), "工程效率", [
        "相比 LSTM 参数更少，训练和推理更轻",
        "相比 Transformer 复杂度更低，先建立稳健 baseline",
        "A100 上 8 epoch 约 4h10m，显存约 22.9GB",
    ], ORANGE)
    add_card(slide, Inches(6.8), Inches(3.85), Inches(5.75), Inches(2.15), "实验结果支撑", [
        "验证 MAE 在 epoch 7 达到 0.006515",
        "短窗口验证 accuracy 达到 94.94%",
        "四天连续闭环仍有 92.41% accuracy",
    ], GREEN)
    add_footer(slide, 3)

    # 4. Data pipeline
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "数据对齐与特征构造", "目标是让每个 110ms 网格行同时表达当前 sensor、网格内 action 和 IO 连续状态。")
    y = Inches(2.0)
    add_process_box(slide, Inches(0.55), y, Inches(1.85), Inches(1.2), "raw sources", "APC\nlivedata\nlog action", LIGHT_BLUE)
    add_arrow(slide, Inches(2.55), y + Inches(0.42), Inches(0.55), Inches(0.28))
    add_process_box(slide, Inches(3.25), y, Inches(2.05), Inches(1.2), "time parse", "毫秒左补零\n按绝对时间排序", LIGHT_TEAL)
    add_arrow(slide, Inches(5.45), y + Inches(0.42), Inches(0.55), Inches(0.28))
    add_process_box(slide, Inches(6.15), y, Inches(2.05), Inches(1.2), "110ms grid", "只输出 livedata\n覆盖片段", LIGHT_BLUE)
    add_arrow(slide, Inches(8.35), y + Inches(0.42), Inches(0.55), Inches(0.28))
    add_process_box(slide, Inches(9.05), y, Inches(1.95), Inches(1.2), "state pass", "覆盖外 action\n先更新 state", LIGHT_ORANGE)
    add_arrow(slide, Inches(11.15), y + Inches(0.42), Inches(0.55), Inches(0.28))
    add_process_box(slide, Inches(11.85), y, Inches(1.0), Inches(1.2), "cache", "numpy\narrays", LIGHT_TEAL)
    add_card(slide, Inches(0.9), Inches(4.05), Inches(3.65), Inches(1.65), "输入特征 401 维", [
        "sensor: 150",
        "evt: 122",
        "state: 122",
        "source one-hot: 4",
        "time features: 3",
    ], BLUE)
    add_card(slide, Inches(4.85), Inches(4.05), Inches(3.65), Inches(1.65), "mask 策略", [
        "mask 不作为模型输入",
        "input_mask 单独保存",
        "用于 loss/accuracy 有效点筛选",
    ], TEAL)
    add_card(slide, Inches(8.8), Inches(4.05), Inches(3.65), Inches(1.65), "异常处理", [
        "仅接受固定 mask 模板的 raw sensor 行",
        "20260416 异常 APC rows 被排除",
        "5/8-5/11 rejected livedata 为 0",
    ], ORANGE)
    add_footer(slide, 4)

    # 5. Architecture diagram
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "GRU 模型框架图", "warmup 阶段使用真实历史初始化 hidden；rollout 阶段将预测 sensor 回灌到下一步输入。")
    y = Inches(1.95)
    add_process_box(slide, Inches(0.55), y, Inches(2.0), Inches(1.25), "Feature x_t", "401 维\n无 mask 输入", LIGHT_BLUE)
    add_arrow(slide, Inches(2.68), y + Inches(0.46), Inches(0.55), Inches(0.28), BLUE)
    add_process_box(slide, Inches(3.35), y, Inches(2.0), Inches(1.25), "GRU layer 1", "hidden=256", LIGHT_TEAL)
    add_arrow(slide, Inches(5.48), y + Inches(0.46), Inches(0.55), Inches(0.28), BLUE)
    add_process_box(slide, Inches(6.15), y, Inches(2.0), Inches(1.25), "GRU layer 2", "dropout=0.1\nhidden=256", LIGHT_TEAL)
    add_arrow(slide, Inches(8.28), y + Inches(0.46), Inches(0.55), Inches(0.28), BLUE)
    add_process_box(slide, Inches(8.95), y, Inches(1.75), Inches(1.25), "Head", "LayerNorm\nLinear", LIGHT_ORANGE)
    add_arrow(slide, Inches(10.82), y + Inches(0.46), Inches(0.55), Inches(0.28), BLUE)
    add_process_box(slide, Inches(11.5), y, Inches(1.35), Inches(1.25), "ŷ_t", "150 sensors\nclamp [0,1]", LIGHT_BLUE)
    fb = slide.shapes.add_shape(MSO_SHAPE.U_TURN_ARROW, Inches(3.15), Inches(3.75), Inches(7.8), Inches(1.15))
    fb.fill.solid()
    fb.fill.fore_color.rgb = RGBColor(230, 238, 248)
    fb.line.color.rgb = RGBColor(190, 205, 225)
    add_textbox(slide, Inches(4.05), Inches(4.15), Inches(5.95), Inches(0.32),
                "rollout: 用预测 sensor 替换下一行 sensor 输入，action/state/time 特征沿 aligned 序列提供", 11, BLUE,
                True, align=PP_ALIGN.CENTER)
    add_card(slide, Inches(0.8), Inches(5.35), Inches(3.75), Inches(1.15), "训练目标", [
        "预测归一化 sensor 绝对值",
        "masked weighted Huber",
    ], TEAL)
    add_card(slide, Inches(4.85), Inches(5.35), Inches(3.75), Inches(1.15), "窗口设置", [
        "warmup=256 steps",
        "rollout=64 steps",
    ], BLUE)
    add_card(slide, Inches(8.9), Inches(5.35), Inches(3.45), Inches(1.15), "部署含义", [
        "保存 hidden 即可连续推理",
        "支持跨 segment / 跨天继承",
    ], ORANGE)
    add_footer(slide, 5)

    # 6. Training setup
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "训练配置摘要", "复刻旧目录较优训练策略，但输入去掉 mask 位，target_mask 只用于指标和 loss。")
    add_metric(slide, Inches(0.8), Inches(1.55), Inches(2.1), Inches(1.2), "input_dim", "401", "no mask", BLUE)
    add_metric(slide, Inches(3.2), Inches(1.55), Inches(2.1), Inches(1.2), "output_dim", "150", "sensor", TEAL)
    add_metric(slide, Inches(5.6), Inches(1.55), Inches(2.1), Inches(1.2), "hidden", "256", "2 layers", ORANGE)
    add_metric(slide, Inches(8.0), Inches(1.55), Inches(2.1), Inches(1.2), "warmup", "256", "28.16s", GREEN)
    add_metric(slide, Inches(10.4), Inches(1.55), Inches(2.1), Inches(1.2), "rollout", "64", "7.04s", RED)
    add_table(slide, Inches(0.8), Inches(3.25), Inches(11.75), Inches(2.3),
              ["参数", "取值", "说明"],
              [
                  ["optimizer", "AdamW", "lr=3e-4, weight_decay=1e-4"],
                  ["batch_size", "64", "A100 bf16 mixed precision"],
                  ["train_stride / val_stride", "8 / 64", "滑动窗口采样"],
                  ["loss", "masked weighted Huber", "active/change continuous 加权"],
                  ["best checkpoint", "epoch 7", "val_mae=0.006515"],
              ],
              font_size=10)
    add_footer(slide, 6)

    # 7. Training convergence
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "训练收敛与停止依据", "训练集继续下降，但验证集在 epoch 7 后变差，因此采用 epoch 7 的 best.pt。")
    epochs = ["1", "2", "3", "4", "5", "6", "7", "8"]
    val_mae = [0.008109, 0.007534, 0.007251, 0.006939, 0.006704, 0.006630, 0.006515, 0.006574]
    train_mae = [0.008720, 0.005112, 0.004537, 0.004193, 0.003980, 0.003814, 0.003701, 0.003649]
    add_bar_chart(slide, Inches(0.75), Inches(1.45), Inches(5.75), Inches(3.9), epochs, val_mae, 0.009,
                  "val_mae by epoch", BLUE)
    add_bar_chart(slide, Inches(6.85), Inches(1.45), Inches(5.75), Inches(3.9), epochs, train_mae, 0.009,
                  "train_mae by epoch", TEAL)
    add_card(slide, Inches(0.95), Inches(5.65), Inches(5.4), Inches(1.1), "停止依据", [
        "epoch 7: val_mae=0.006515，为当前 best",
        "epoch 8: val_mae 回升到 0.006574",
    ], ORANGE)
    add_card(slide, Inches(6.95), Inches(5.65), Inches(5.4), Inches(1.1), "保留模型", [
        "runs/gru_0430_replicate_nomask/checkpoints/best.pt",
        "对应 epoch 7，后续评估均使用该 checkpoint",
    ], GREEN)
    add_footer(slide, 7)

    # 8. Millisecond fix
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "livedata 毫秒修复与验证集 cache", "问题只影响 5/8-5/11 验证侧，因此只重建 val-only cache，训练 cache 沿用原结果。")
    add_table(slide, Inches(0.75), Inches(1.45), Inches(5.5), Inches(1.55),
              ["raw time", "fixed parse"],
              [
                  ["2026/5/8 21:05:04:7", "2026-05-08 21:05:04.007"],
                  ["2026/5/9 11:30:25:87", "2026-05-09 11:30:25.087"],
                  ["2026/5/9 00:00:16.893", "2026-05-09 00:00:16.893"],
              ],
              font_size=9, header_fill=TEAL)
    add_metric(slide, Inches(6.75), Inches(1.45), Inches(1.75), Inches(1.15), "raw rows", "172k", "5/8-5/9", BLUE)
    add_metric(slide, Inches(8.75), Inches(1.45), Inches(1.75), Inches(1.15), "1/2 digit ms", "17.4k", "需左补零", ORANGE)
    add_metric(slide, Inches(10.75), Inches(1.45), Inches(1.75), Inches(1.15), "倒序", "0", "修复后", GREEN)
    add_table(slide, Inches(0.75), Inches(3.45), Inches(11.75), Inches(2.0),
              ["date", "aligned rows", "hidden_actions", "livedata_grid", "apc_grid"],
              [
                  ["20260508", "489,265", "29", "295,477", "193,788"],
                  ["20260509", "784,752", "0", "622,096", "162,656"],
                  ["20260510", "360,552", "4,064", "355,452", "5,100"],
                  ["20260511", "501,443", "0", "473,471", "27,972"],
              ],
              font_size=9)
    add_card(slide, Inches(0.85), Inches(5.8), Inches(11.55), Inches(0.8), "val-only cache 校验", [
        "rows=2,136,012; features=(2136012, 401); valid_targets=2,135,991; masks_in_features=False; ts negative=0; duplicate=0",
    ], BLUE)
    add_footer(slide, 8)

    # 9. Evaluation modes
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "全天自回归评估方式", "评估从短窗口扩展到全天和四天连续，重点验证误差在长链中的累积。")
    add_card(slide, Inches(0.75), Inches(1.55), Inches(3.8), Inches(4.65), "短窗口验证", [
        "warmup=256, rollout=64",
        "窗口内部闭环",
        "旧验证 accuracy=94.94%",
        "用于确认模型基本预测能力",
    ], BLUE)
    add_card(slide, Inches(4.75), Inches(1.55), Inches(3.8), Inches(4.65), "每天单独全天", [
        "每一天开头 warmup 256 步",
        "日内 segment/gap 不重置",
        "跨天重置",
        "用于观察单日长链稳定性",
    ], TEAL)
    add_card(slide, Inches(8.75), Inches(1.55), Inches(3.8), Inches(4.65), "四天连续全天", [
        "只在 5/8 开头 warmup 一次",
        "5/9-5/11 不重新 warmup",
        "跨天继承 hidden 和预测 sensor",
        "用于模拟更严格的连续运行",
    ], ORANGE)
    add_footer(slide, 9)

    # 10. Daily results
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "每天单独全天自回归结果", "总体 accuracy=93.12%，continuous sensor 是主要误差来源，binary sensor 保持较高准确率。")
    add_table(slide, Inches(0.65), Inches(1.45), Inches(12.05), Inches(2.45),
              ["date", "steps", "accuracy", "cont_acc", "bin_acc", "MAE"],
              [
                  ["20260508", "489,009", "91.70%", "88.81%", "99.49%", "0.013216"],
                  ["20260509", "784,496", "92.77%", "90.22%", "99.68%", "0.008941"],
                  ["20260510", "360,296", "95.20%", "93.43%", "99.97%", "0.003853"],
                  ["20260511", "501,187", "93.57%", "91.23%", "99.90%", "0.009882"],
                  ["overall", "2,134,988", "93.12%", "90.67%", "99.74%", "0.009283"],
              ],
              font_size=10)
    add_bar_chart(slide, Inches(0.75), Inches(4.35), Inches(5.6), Inches(2.1),
                  ["5/8", "5/9", "5/10", "5/11"], [0.917, 0.928, 0.952, 0.936], 1.0,
                  "daily accuracy", GREEN)
    add_bar_chart(slide, Inches(6.85), Inches(4.35), Inches(5.6), Inches(2.1),
                  ["5/8", "5/9", "5/10", "5/11"], [0.0132, 0.0089, 0.0039, 0.0099], 0.016,
                  "daily MAE", ORANGE)
    add_footer(slide, 10)

    # 11. Continuous results
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "四天连续自回归结果", "四天只 warmup 一次，跨天不重置；结果低于每天单独评估，体现长链误差积累。")
    add_metric(slide, Inches(0.8), Inches(1.5), Inches(2.2), Inches(1.25), "accuracy", "92.41%", "4-day continuous", BLUE)
    add_metric(slide, Inches(3.35), Inches(1.5), Inches(2.2), Inches(1.25), "cont_acc", "89.83%", "continuous sensors", TEAL)
    add_metric(slide, Inches(5.9), Inches(1.5), Inches(2.2), Inches(1.25), "bin_acc", "99.37%", "binary sensors", GREEN)
    add_metric(slide, Inches(8.45), Inches(1.5), Inches(2.2), Inches(1.25), "MAE", "0.01129", "overall", ORANGE)
    add_metric(slide, Inches(11.0), Inches(1.5), Inches(1.65), Inches(1.25), "steps", "2.136M", "one chain", RED)
    add_table(slide, Inches(1.05), Inches(3.35), Inches(11.15), Inches(1.3),
              ["评估方式", "reset 策略", "accuracy", "cont_acc", "bin_acc", "MAE"],
              [
                  ["每天单独全天", "每天开头 warmup；日内不重置", "93.12%", "90.67%", "99.74%", "0.009283"],
                  ["四天连续全天", "5/8 warmup 一次；跨天不重置", "92.41%", "89.83%", "99.37%", "0.011290"],
              ],
              font_size=10)
    add_card(slide, Inches(1.05), Inches(5.15), Inches(5.25), Inches(1.15), "观察", [
        "连续评估比每日评估低 0.72 个百分点",
        "MAE 从 0.009283 上升到 0.011290",
    ], ORANGE)
    add_card(slide, Inches(6.7), Inches(5.15), Inches(5.5), Inches(1.15), "解释", [
        "跨天不重置会积累 sensor 预测误差",
        "连续 sensor 比 binary sensor 更敏感",
    ], TEAL)
    add_footer(slide, 11)

    # 12. Conclusions
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "结论与下一步", "GRU baseline 已可用于全天闭环验证；下一阶段应围绕长链漂移和 continuous sensor 误差优化。")
    add_card(slide, Inches(0.75), Inches(1.45), Inches(3.8), Inches(4.7), "已确认", [
        "livedata 毫秒解析问题已修复",
        "5/8-5/11 aligned/cache 时间顺序无异常",
        "val-only cache 输入维度为 401，无 mask 输入",
        "best GRU checkpoint 可稳定完成全天闭环",
    ], GREEN)
    add_card(slide, Inches(4.75), Inches(1.45), Inches(3.8), Inches(4.7), "当前结论", [
        "短窗口 accuracy 94.94%",
        "每天全天 accuracy 93.12%",
        "四天连续 accuracy 92.41%",
        "跨天连续运行存在可见误差积累",
    ], BLUE)
    add_card(slide, Inches(8.75), Inches(1.45), Inches(3.8), Inches(4.7), "建议", [
        "优先分析 continuous sensor 长链漂移",
        "按动作后变化区间做误差分桶",
        "尝试调整 active/change loss 权重",
        "后续再与更复杂序列模型做对比",
    ], ORANGE)
    add_footer(slide, 12)

    return prs


def main() -> None:
    prs = build_presentation()
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(OUT_PATH)


if __name__ == "__main__":
    main()
